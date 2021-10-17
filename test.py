

'''
    Importing required modules
'''
from pptx import Presentation
from pptx.util import Inches

'''
    Creating pptx object
'''
pr1 = Presentation()

# Slide contents to insert
slide_content = {
    'slide_1':{
        'title':'Image Title 1',
        'subtitle':'Image Subtitle 1',
        'image':'output1.jpg'
    },
    'slide_2':{
        'title':'Image Title 2',
        'subtitle':'Image Subtitle 2',
        'image':'output2.jpg'
    },
    'slide_3':{
        'title':'Image Title 3',
        'subtitle':'Image Subtitle 3',
        'image':'output3.jpg'
    },
    'slide_4':{
        'title':'Image Title 4',
        'subtitle':'Image Subtitle 4',
        'image':'output4.jpg'
    },
    'slide_5':{
        'title':'Image Title 5',
        'subtitle':'Image Subtitle 5',
        'image':'output5.jpg'
    }
}

'''
     Creating slide with slide content
'''
for x in slide_content:
    slide_register = pr1.slide_layouts[1]
    slide_1 = pr1.slides.add_slide(slide_register)
    title = slide_1.shapes.title
    subtitle = slide_1.placeholders[1]
    title.text = slide_content[x]['title']
    subtitle.text= slide_content[x]['subtitle']
    img1=(slide_content[x]['image'])
    from_left = Inches(1)
    from_top = Inches(3)
    add_picture = slide_1.shapes.add_picture(img1,from_left,from_top)

# saving the pptx file
try:
    pr1.save("Python_ppt.pptx")
    print('PPT FILE IS READY')
except:
    print('Opps! Something went wrong.')